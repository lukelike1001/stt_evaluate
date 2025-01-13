import torchaudio
from torchaudio.functional import resample
from transformers import AutoModelForSpeechSeq2Seq, AutoProcessor
from jiwer import wer
from pathlib import Path
import einops
import requests
import signal
import os
import subprocess
import json
import xlsxwriter
from pptx import Presentation
from pptx.util import Pt

system_prompt = (
        "**GPT Agent Prompt**\n\n"
        "You are tasked with analyzing a transcript and producing a 2D JSON structure that represents an F2T2TEA table. "
        "The table should contain seven fields: Enemy, Find, Fix, Track, Target, Engage, and Assess. The JSON output must "
        "be a list of dictionaries where each dictionary represents a row in the table, with keys corresponding to the column "
        "names and values corresponding to the data extracted from the transcript.\n\n"
        "### Column Definitions:\n"
        "1. **Enemy**: \n"
        "   - The name or type of enemy unit mentioned in the transcript, such as Enemy1, Enemy2, etc.\n"
        "   - Each enemy mentioned in the transcript should have a unique dictionary entry.\n\n"
        "2. **Find**:\n"
        "   - The unit responsible for the initial detection or identification of the enemy.\n"
        "   - Look for phrases like \"can find,\" \"for Find,\" or \"what do you have for that?\" and match them to the enemy unit.\n\n"
        "3. **Fix**:\n"
        "   - The unit responsible for maintaining surveillance or locking onto the enemy once detected.\n"
        "   - Look for statements that refer to \"Fix\" or mention a similar responsibility.\n\n"
        "4. **Track**:\n"
        "   - The unit responsible for ongoing monitoring of the enemy’s movements.\n"
        "   - Look for phrases like \"track the enemy\" or \"can track.\"\n\n"
        "5. **Target**:\n"
        "   - The unit tasked with planning or preparing to attack the enemy.\n"
        "   - Look for phrases like \"target Enemy1\" or \"you’ll target.\"\n\n"
        "6. **Engage**:\n"
        "   - The unit actively attacking the enemy.\n"
        "   - Look for phrases like \"engage Enemy1\" or \"will engage.\"\n\n"
        "7. **Assess**:\n"
        "   - The unit responsible for evaluating the result of the attack or confirming destruction.\n"
        "   - Look for mentions of \"assess\" or evaluation after an attack, including visual confirmation by aircraft or other means.\n\n"
        "### Output Requirements:\n"
        "- The output must be a valid 2D JSON structure in the form of a list of dictionaries.\n"
        "- Each dictionary should have keys: Enemy, Find, Fix, Track, Target, Engage, Assess.\n"
        "- Values should be filled based on the transcript. If no value is explicitly assigned for a particular column, leave it as an empty string (\"\").\n"
        "- The JSON output must be provided as a list of dictionaries without any additional wording, explanations, or codeblocks.\n\n"
        "### Example Output Format:\n"
        "[\n"
        "    {\"Enemy\": \"Enemy1\", \"Find\": \"UnitA\", \"Fix\": \"UnitB\", \"Track\": \"UnitC\", \"Target\": \"UnitD\", \"Engage\": \"UnitE\", \"Assess\": \"UnitF\"},\n"
        "    {\"Enemy\": \"Enemy2\", \"Find\": \"UnitG\", \"Fix\": \"UnitH\", \"Track\": \"UnitI\", \"Target\": \"UnitJ\", \"Engage\": \"UnitK\", \"Assess\": \"UnitL\"},\n"
        "    {\"Enemy\": \"Enemy3\", \"Find\": \"UnitM\", \"Fix\": \"UnitN\", \"Track\": \"UnitO\", \"Target\": \"UnitP\", \"Engage\": \"UnitQ\", \"Assess\": \"UnitR\"},\n"
        "    {\"Enemy\": \"Enemy4\", \"Find\": \"UnitS\", \"Fix\": \"UnitT\", \"Track\": \"UnitU\", \"Target\": \"UnitV\", \"Engage\": \"UnitW\", \"Assess\": \"UnitX\"}\n"
        "]\n\n"
        "MAKE SURE TO ONLY USE CODE NAMES EVEN IF THEY ARE REFERENCED AS SOMETHING ELSE WIHTIN THE TRANSCRIPT. Here is a dictionary that you can use to determine the code names:"
        + """**Aircraft Units:**

        1. **EA-18G (3 units)**
        - Unit: VAQ-135
        - Codename: Hawk

        2. **B-1B (2 units)**
        - Unit: 34 BS
        - Codename: Eagle

        3. **F-22A (4 units)**
        - Unit: 7 FS
        - Codename: Thunderwave

        4. **F-35 (4 units)**
        - Unit: VFA-147
        - Codename: Knight

        5. **EA-18G (2 units)**
        - Unit: VAQ-141
        - Codename: Wolverine

        6. **FA-18E/F (8 units)**
        - Unit: CVW-5
        - Codename: Lightning

        7. **E-2D (1 unit)**
        - Unit: VAW-125
        - Codename: Bluejay

        8. **E-7A (1 unit)**
        - Unit: 2 SQN RAAF
        - Codename: Rook

        9. **KC-46 (2 units)**
        - Unit: TTF
        - Codename: Dragonfly

        10. **KC-135 (2 units)**
            - Unit: TTF
            - Codename: Hornet

        **Enemy Units:**

        1. **Su-35s**
        - Codename: Salmon

        2. **J-11s**
        - Codename: Swallow

        3. **JH-7s**
        - Codename: Sparrow

        4. **CH-SA-21**
        - Codename: Swordfish
        """
    ) 

class TranscriptProcessor:
    def __init__(self):
        pass

    def generate_completion(self, prompt):
        """
        Sends a request to the Ollama API to generate a text completion.

        Args:
            prompt (str): The input prompt for the model.

        Returns:
            str: The response from the Ollama API.
        """
        url = "http://localhost:11434/api/generate"

        payload = {
            "model": "llama3.1:70b",
            "prompt": prompt,
            "stream": False,
            "options": {
                "temperature": 0.7  # Default temperature
            }
        }

        try:
            response = requests.post(url, json=payload, stream=False)
            response.raise_for_status()

            # Return a single response object
            return response.json().get("response", "No response received.")

        except requests.exceptions.RequestException as e:
            print(f"Error: {e}")
            return None

    def json_to_excel(self, json_data, excel_file):
        """
        Converts a JSON structure to an Excel file.

        Args:
            json_data (list of dict): The JSON data to write.
            excel_file (str): Path to the Excel file to create.
        """
        workbook = xlsxwriter.Workbook(excel_file)
        worksheet = workbook.add_worksheet()

        headers = list(json_data[0].keys()) if json_data else []

        for col, header in enumerate(headers):
            worksheet.write(0, col, header)

        for row_num, record in enumerate(json_data, start=1):
            for col_num, header in enumerate(headers):
                worksheet.write(row_num, col_num, record.get(header, ''))

        workbook.close()

    def json_to_slides(self, json_data, pptx_file="presentation.pptx"):
        """
        Converts each entry in the JSON data into a slide, with detailed content generated
        via LLM calls for each entry, and saves it as a .pptx file.

        Args:
            json_data (list of dict): The JSON data to process.
            pptx_file (str): The path to the .pptx file to create.

        Returns:
            None
        """
        presentation = Presentation()

        for idx, entry in enumerate(json_data, start=1):
            prompt = (
                f"You are creating content for a slide presentation. This is slide {idx}. "
                f"Based on the following table entry, generate a concise, professional, and "
                f"well-structured summary suitable for a slide. Include details about the enemy unit "
                f"and the assigned actions by various units. Format the output as text suitable for a slide.\n\n"
                f"Enemy: {entry['Enemy']}\n"
                f"Find: {entry['Find']}\n"
                f"Fix: {entry['Fix']}\n"
                f"Track: {entry['Track']}\n"
                f"Target: {entry['Target']}\n"
                f"Engage: {entry['Engage']}\n"
                f"Assess: {entry['Assess']}\n"
                f"\nGenerate a clear and professional description for this slide."
            )

            slide_content = self.generate_completion(prompt)

            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            text_box = slide.shapes.add_textbox(Pt(50), Pt(50), Pt(600), Pt(400))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            if slide_content:
                text_frame.text = f"Slide {idx}: {entry['Enemy']}"
                p = text_frame.add_paragraph()
                p.text = slide_content
            else:
                text_frame.text = f"Slide {idx}: {entry['Enemy']}"
                p = text_frame.add_paragraph()
                p.text = "Failed to generate content."

        presentation.save(pptx_file)
        print(f"Presentation saved as '{pptx_file}'")

    def stop_llm_process(self):
        """
        Sends a signal to stop the LLM process running on the specified port.
        """
        try:
            subprocess.run(["pkill", "-f", "ollama-server"], check=True)
            print("LLM process stopped successfully.")
        except subprocess.CalledProcessError:
            print("LLM process is not running or could not be stopped.")

    def evaluate_whisper_base(self, input_audio, input_reference):
        """
        Evaluate Whisper base model on an audio file and reference transcript.

        Args:
            input_audio (str): Path to the input audio file.
            input_reference (str): Reference transcript.

        Returns:
            tuple: Transcription and Word Error Rate (WER).
        """
        audio, sr = torchaudio.load(input_audio)

        if sr != 16000:
            audio = resample(audio, sr, 16000)
            print(f"Resampled audio to 16kHz")

        model = AutoModelForSpeechSeq2Seq.from_pretrained("openai/whisper-base.en")
        processor = AutoProcessor.from_pretrained("openai/whisper-base.en")

        print("Preprocessing audio...")
        input_features = processor.feature_extractor(audio.squeeze().numpy(), sampling_rate=16000, return_tensors="pt").input_features

        print("Transcribing audio...")
        predicted_ids = model.generate(input_features)
        transcription = processor.tokenizer.batch_decode(predicted_ids, skip_special_tokens=True)[0]
        print("Transcription:", transcription)

        reference_transcript = input_reference.strip()
        error_rate = wer(reference_transcript, transcription)
        print(f"Word Error Rate (WER): {error_rate:.2%}")

        return transcription, error_rate

    def process_transcription_pipeline(self, audio_files, reference_transcripts, excel_file="output.xlsx", pptx_file="presentation.pptx"):
        """
        Processes audio files through transcription, generates JSON, and creates Excel and PowerPoint slides.

        Args:
            audio_files (list of str): Paths to audio files.
            reference_transcripts (list of str): Corresponding reference transcripts.
            excel_file (str): Output Excel file path.
            pptx_file (str): Output PowerPoint file path.

        Returns:
            None
        """
        json_data = []

        # Step 1: Transcribe audio and calculate WER
        for idx, (audio_file, reference) in enumerate(zip(audio_files, reference_transcripts), start=1):
            print(f"Processing audio file {idx}/{len(audio_files)}: {audio_file}")
            transcription, wer_score = self.evaluate_whisper_base(audio_file, reference)
            print(f"Transcription: {transcription}, WER: {wer_score:.2%}")

            combined_prompt = f"{system_prompt}\n\n----\nTranscription to be turned into json:{transcription}"
            
            json_response = self.generate_completion(combined_prompt)
        # Step 2: Save JSON to Excel
        self.json_to_excel(json_data, excel_file)
        print(f"Excel file saved to {excel_file}")

        # Step 3: Generate slides from JSON
        self.json_to_slides(json_data, pptx_file)
        print(f"PowerPoint presentation saved to {pptx_file}")

if __name__ == "__main__":
    processor = TranscriptProcessor()

    # Example usage of the process_transcription_pipeline
    audio_files = [
        "path/to/audio1.wav",
    ]
    reference_transcripts = [
        "Expected transcript for audio 1.",
    ]

    processor.process_transcription_pipeline(audio_files, reference_transcripts)
