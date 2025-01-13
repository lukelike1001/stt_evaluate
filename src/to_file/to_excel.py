import requests
import signal
import os
import subprocess
import json
import xlsxwriter

def generate_completion(prompt):
    """
    Sends a request to the Ollama API to generate a text completion.

    Args:
        prompt (str): The input prompt for the model.

    Returns:
        str: The response from the Ollama API.
    """
    url = "http://localhost:11434/api/generate"

    payload = {
        "model": "tinyllama",
        "prompt": prompt,
        "stream": False,
        "options": {
            "temperature": 0  # Default temperature
        }
    }

    try:
        response = requests.post(url, json=payload, stream=False)
        response.raise_for_status()

        # Return a single response object
        return response.json().get("response", "No response received.")

    except requests.exceptions.RequestException as e:
        print(f"Error: {e}")
        return None

def json_to_excel(json_data, excel_file):
    """
    Converts a JSON structure to an Excel file.

    Args:
        json_data (list of dict): The JSON data to write.
        excel_file (str): Path to the Excel file to create.
    """
    # Create a workbook and add a worksheet
    workbook = xlsxwriter.Workbook(excel_file)
    worksheet = workbook.add_worksheet()

    # Extract keys and preserve their order from the first record
    headers = list(json_data[0].keys()) if json_data else []

    # Write headers to the first row
    for col, header in enumerate(headers):
        worksheet.write(0, col, header)

    # Write data rows
    for row_num, record in enumerate(json_data, start=1):
        for col_num, header in enumerate(headers):
            worksheet.write(row_num, col_num, record.get(header, ''))

    # Close the workbook
    workbook.close()

def stop_llm_process():
    """
    Sends a signal to stop the LLM process running on the specified port.
    """
    try:
        # Attempt to stop the server by sending a termination signal
        subprocess.run(["pkill", "-f", "ollama-server"], check=True)
        print("LLM process stopped successfully.")
    except subprocess.CalledProcessError:
        print("LLM process is not running or could not be stopped.")
        
from pptx import Presentation
from pptx.util import Pt

def json_to_slides(json_data, pptx_file="presentation.pptx"):
    """
    Converts each entry in the JSON data into a slide, with detailed content generated
    via LLM calls for each entry, and saves it as a .pptx file.

    Args:
        json_data (list of dict): The JSON data to process.
        pptx_file (str): The path to the .pptx file to create.

    Returns:
        None
    """
    # Create a new PowerPoint presentation
    presentation = Presentation()

    for idx, entry in enumerate(json_data, start=1):
        # Generate content for the slide using LLM
        prompt = (
            f"You are creating content for a slide presentation. This is slide {idx}. "
            f"Based on the following table entry, generate a concise, professional, and "
            f"well-structured summary suitable for a slide. Include details about the enemy unit "
            f"and the assigned actions by various units. Format the output as text suitable for a slide.\n\n"
            f"Enemy: {entry['Enemy']}\n"
            f"Find: {entry['Find']}\n"
            f"Fix: {entry['Fix']}\n"
            f"Track: {entry['Track']}\n"
            f"Target: {entry['Target']}\n"
            f"Engage: {entry['Engage']}\n"
            f"Assess: {entry['Assess']}\n"
            f"\nGenerate a clear and professional description for this slide."
        )

        slide_content = generate_completion(prompt)

        # Add a slide for each entry
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Use blank slide layout
        text_box = slide.shapes.add_textbox(Pt(50), Pt(50), Pt(600), Pt(400))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        if slide_content:
            text_frame.text = f"Slide {idx}: {entry['Enemy']}"
            p = text_frame.add_paragraph()
            p.text = slide_content
        else:
            text_frame.text = f"Slide {idx}: {entry['Enemy']}"
            p = text_frame.add_paragraph()
            p.text = "Failed to generate content."

    # Save the presentation to a file
    presentation.save(pptx_file)
    print(f"Presentation saved as '{pptx_file}'")


if __name__ == "__main__":
    system_prompt = (
        "**GPT Agent Prompt**\n\n"
        "You are tasked with analyzing a transcript and producing a 2D JSON structure that represents an F2T2TEA table. "
        "The table should contain seven fields: Enemy, Find, Fix, Track, Target, Engage, and Assess. The JSON output must "
        "be a list of dictionaries where each dictionary represents a row in the table, with keys corresponding to the column "
        "names and values corresponding to the data extracted from the transcript.\n\n"
        "### Column Definitions:\n"
        "1. **Enemy**: \n"
        "   - The name or type of enemy unit mentioned in the transcript, such as Enemy1, Enemy2, etc.\n"
        "   - Each enemy mentioned in the transcript should have a unique dictionary entry.\n\n"
        "2. **Find**:\n"
        "   - The unit responsible for the initial detection or identification of the enemy.\n"
        "   - Look for phrases like \"can find,\" \"for Find,\" or \"what do you have for that?\" and match them to the enemy unit.\n\n"
        "3. **Fix**:\n"
        "   - The unit responsible for maintaining surveillance or locking onto the enemy once detected.\n"
        "   - Look for statements that refer to \"Fix\" or mention a similar responsibility.\n\n"
        "4. **Track**:\n"
        "   - The unit responsible for ongoing monitoring of the enemy’s movements.\n"
        "   - Look for phrases like \"track the enemy\" or \"can track.\"\n\n"
        "5. **Target**:\n"
        "   - The unit tasked with planning or preparing to attack the enemy.\n"
        "   - Look for phrases like \"target Enemy1\" or \"you’ll target.\"\n\n"
        "6. **Engage**:\n"
        "   - The unit actively attacking the enemy.\n"
        "   - Look for phrases like \"engage Enemy1\" or \"will engage.\"\n\n"
        "7. **Assess**:\n"
        "   - The unit responsible for evaluating the result of the attack or confirming destruction.\n"
        "   - Look for mentions of \"assess\" or evaluation after an attack, including visual confirmation by aircraft or other means.\n\n"
        "### Output Requirements:\n"
        "- The output must be a valid 2D JSON structure in the form of a list of dictionaries.\n"
        "- Each dictionary should have keys: Enemy, Find, Fix, Track, Target, Engage, Assess.\n"
        "- Values should be filled based on the transcript. If no value is explicitly assigned for a particular column, leave it as an empty string (\"\").\n"
        "- The JSON output must be provided as a list of dictionaries without any additional wording, explanations, or codeblocks.\n\n"
        "### Example Output Format:\n"
        "[\n"
        "    {\"Enemy\": \"Enemy1\", \"Find\": \"UnitA\", \"Fix\": \"UnitB\", \"Track\": \"UnitC\", \"Target\": \"UnitD\", \"Engage\": \"UnitE\", \"Assess\": \"UnitF\"},\n"
        "    {\"Enemy\": \"Enemy2\", \"Find\": \"UnitG\", \"Fix\": \"UnitH\", \"Track\": \"UnitI\", \"Target\": \"UnitJ\", \"Engage\": \"UnitK\", \"Assess\": \"UnitL\"},\n"
        "    {\"Enemy\": \"Enemy3\", \"Find\": \"UnitM\", \"Fix\": \"UnitN\", \"Track\": \"UnitO\", \"Target\": \"UnitP\", \"Engage\": \"UnitQ\", \"Assess\": \"UnitR\"},\n"
        "    {\"Enemy\": \"Enemy4\", \"Find\": \"UnitS\", \"Fix\": \"UnitT\", \"Track\": \"UnitU\", \"Target\": \"UnitV\", \"Engage\": \"UnitW\", \"Assess\": \"UnitX\"}\n"
        "]\n\n"
        "MAKE SURE TO ONLY USE UNIT NAMES"
    )

    user_prompt = input("Enter the transcript for analysis: ")

    # Combine system prompt and user prompt
    combined_prompt = f"{system_prompt}\n\n{user_prompt}"

    stop_signal = input("Do you want to stop the LLM process before proceeding? (yes/no): ").strip().lower()
    if stop_signal == "yes":
        stop_llm_process()

    json_response = generate_completion(combined_prompt)

    if json_response:
        try:
            print(json_response)
            parsed_json = json.loads(json_response)
            excel_file = "F2T2TEA_table.xlsx"
            json_to_excel(parsed_json, excel_file)
            print(f"Excel file '{excel_file}' created successfully.")
        except json.JSONDecodeError:
            print("Failed to parse JSON response. Ensure the output is valid JSON.")
    else:
        print("No response received from the LLM.")

    json_data = [
        {"Enemy": "Su-35", "Find": "VAQ-135", "Fix": "VAQ-135", "Track": "", "Target": "VFA-147", "Engage": "VFA-147", "Assess": "VFA-147"},
        {"Enemy": "J-11", "Find": "2 SQN RAAF", "Fix": "2 SQN RAAF", "Track": "7FS", "Target": "7FS", "Engage": "7FS", "Assess": "34BS"},
        {"Enemy": "JH-7", "Find": "VAQ-141", "Fix": "VAW-125", "Track": "VAW-125", "Target": "7FS", "Engage": "7FS", "Assess": "34BS"},
        {"Enemy": "CH-SA-21", "Find": "CVW-5", "Fix": "CVW-5", "Track": "CVW-5", "Target": "CVW-5", "Engage": "VFA-147", "Assess": ""},
    ]

    json_to_slides(json_data, pptx_file="F2T2TEA_presentation.pptx")
